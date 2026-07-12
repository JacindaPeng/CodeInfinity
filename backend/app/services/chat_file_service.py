"""对话附件解析：从上传文件中提取文本供大模型参考。"""
from __future__ import annotations

import os
import shutil
import subprocess
import sys
import tempfile
import uuid
from pathlib import Path

from fastapi import HTTPException, UploadFile

from ..config import settings

# 文档类（.doc 为旧版 Word，需先转换为 docx 再解析）
_DOC_EXTS = {".pdf", ".ppt", ".pptx", ".doc", ".docx", ".txt", ".md", ".markdown"}
# 代码 / 数据类（按纯文本读取）
_TEXT_EXTS = {
    ".c", ".h", ".cpp", ".cc", ".java", ".py", ".js", ".ts", ".vue",
    ".html", ".css", ".json", ".xml", ".csv", ".log", ".sql", ".sh", ".bat",
}
# 图片类（暂不做 OCR，返回提示信息）
_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"}

ALLOWED_EXTS = _DOC_EXTS | _TEXT_EXTS | _IMAGE_EXTS
MAX_FILE_BYTES = 10 * 1024 * 1024
MAX_TEXT_CHARS = 50_000


def chat_attach_dir(user_id: int) -> Path:
    return Path(settings.upload_dir).parent / "chat_attachments" / str(user_id)


def save_chat_attachment(user_id: int, raw: bytes, filename: str) -> str:
    """持久化对话附件，返回 file_id。"""
    file_id = uuid.uuid4().hex
    safe_name = Path(filename).name.replace("..", "").replace("/", "_").replace("\\", "_")
    dest_dir = chat_attach_dir(user_id)
    dest_dir.mkdir(parents=True, exist_ok=True)
    dest = dest_dir / f"{file_id}__{safe_name}"
    dest.write_bytes(raw)
    return file_id


def resolve_chat_attachment(user_id: int, file_id: str) -> Path | None:
    base = chat_attach_dir(user_id)
    if not base.is_dir():
        return None
    for p in base.glob(f"{file_id}*"):
        if p.is_file():
            return p
    return None


def _read_plain_text(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")


def _extract_pdf(path: str) -> str:
    import pdfplumber
    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page_no, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                parts.append(f"[第 {page_no} 页]\n{text}")
    return "\n\n".join(parts)


def _extract_ppt(path: str) -> str:
    from pptx import Presentation
    parts: list[str] = []
    prs = Presentation(path)
    for slide_no, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            elif getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    texts.append(" ".join(c.text for c in row.cells))
        text = "\n".join(t for t in texts if t.strip())
        if text.strip():
            parts.append(f"[幻灯片 {slide_no}]\n{text}")
    return "\n\n".join(parts)


def _extract_word(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for tbl in doc.tables:
        for row in tbl.rows:
            parts.append(" ".join(c.text for c in row.cells))
    return "\n".join(parts)


def _soffice_candidates() -> list[str]:
    cmds: list[str] = []
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            cmds.append(found)
    if sys.platform == "win32":
        for p in (
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ):
            if os.path.isfile(p):
                cmds.append(p)
    return cmds


def _convert_doc_to_docx(path: str, out_dir: str) -> str | None:
    """将 .doc 转为 .docx，成功返回 docx 路径。"""
    stem = Path(path).stem
    for soffice in _soffice_candidates():
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "docx", "--outdir", out_dir, path],
                check=True,
                capture_output=True,
                timeout=90,
            )
            candidate = Path(out_dir) / f"{stem}.docx"
            if candidate.exists():
                return str(candidate)
        except (subprocess.SubprocessError, OSError):
            continue

    if sys.platform == "win32":
        try:
            import win32com.client  # type: ignore

            out = str(Path(out_dir) / f"{stem}.docx")
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            word.DisplayAlerts = 0
            doc = word.Documents.Open(os.path.abspath(path), ReadOnly=True)
            doc.SaveAs2(out, FileFormat=16)  # wdFormatXMLDocument
            doc.Close(False)
            word.Quit()
            if os.path.isfile(out):
                return out
        except Exception:
            pass
    return None


def _extract_doc(path: str) -> str:
    tmp_dir = tempfile.mkdtemp(prefix="chat_doc_")
    try:
        docx_path = _convert_doc_to_docx(path, tmp_dir)
        if not docx_path:
            raise HTTPException(
                400,
                "无法解析旧版 .doc 文件。Windows 请确认已安装 Microsoft Word；"
                "或安装 LibreOffice；也可在 Word 中将文件另存为 .docx 后上传。",
            )
        return _extract_word(docx_path)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def extract_file_text(path: str, filename: str) -> tuple[str, str, bool]:
    """返回 (file_type, text, truncated)。"""
    ext = Path(filename).suffix.lower()
    if ext in _IMAGE_EXTS:
        return "image", (
            f"（图片文件 {filename}，当前对话为文本模式，无法自动识别图片内容。"
            "请用文字描述图片，或换用支持视觉理解的大模型。）"
        ), False

    if ext == ".pdf":
        text = _extract_pdf(path)
        file_type = "pdf"
    elif ext in (".ppt", ".pptx"):
        text = _extract_ppt(path)
        file_type = "ppt"
    elif ext == ".docx":
        text = _extract_word(path)
        file_type = "word"
    elif ext == ".doc":
        text = _extract_doc(path)
        file_type = "word"
    elif ext in _DOC_EXTS | _TEXT_EXTS:
        text = _read_plain_text(path)
        file_type = "text" if ext in _TEXT_EXTS else "document"
    else:
        raise HTTPException(400, f"不支持的文件类型: {ext}")

    text = text.strip()
    if not text:
        raise HTTPException(400, f"未能从文件 {filename} 中提取到文本内容")

    truncated = False
    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS] + "\n\n…（内容过长，已截断）"
        truncated = True
    return file_type, text, truncated


async def parse_upload(file: UploadFile, user_id: int) -> dict:
    """解析上传文件并返回结构化结果。"""
    filename = file.filename or "unknown"
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_EXTS:
        raise HTTPException(
            400,
            f"不支持的文件类型: {ext or '（无扩展名）'}，"
            f"支持 PDF/Word/PPT/TXT/代码文件及常见图片",
        )

    raw = await file.read()
    if len(raw) > MAX_FILE_BYTES:
        raise HTTPException(400, f"文件过大，单文件上限 {MAX_FILE_BYTES // 1024 // 1024}MB")
    if not raw:
        raise HTTPException(400, "文件为空")

    file_id = save_chat_attachment(user_id, raw, filename)

    suffix = ext or ".bin"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    try:
        tmp.write(raw)
        tmp.close()
        file_type, text, truncated = extract_file_text(tmp.name, filename)
    finally:
        try:
            os.unlink(tmp.name)
        except OSError:
            pass

    return {
        "file_id": file_id,
        "name": filename,
        "type": file_type,
        "text": text,
        "truncated": truncated,
        "size": len(raw),
    }
