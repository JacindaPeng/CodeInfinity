"""文档解析：PDF / PPT / Word → 文本片段（带页码/位置元数据）。"""
from __future__ import annotations

from pathlib import Path
from typing import Iterator

from langchain_text_splitters import RecursiveCharacterTextSplitter

_splitter = RecursiveCharacterTextSplitter(
    chunk_size=600,
    chunk_overlap=80,
    separators=["\n\n", "\n", "。", "；", " ", ""],
)


def _split(text: str, base_meta: dict) -> list[tuple[str, dict]]:
    chunks = _splitter.split_text(text)
    out = []
    for i, c in enumerate(chunks):
        if not c.strip():
            continue
        meta = dict(base_meta)
        meta["chunk_idx"] = i
        out.append((c, meta))
    return out


def parse_pdf(path: str, base_meta: dict) -> list[tuple[str, dict]]:
    import pdfplumber
    out: list[tuple[str, dict]] = []
    with pdfplumber.open(path) as pdf:
        for page_no, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if not text.strip():
                continue
            m = dict(base_meta); m["page"] = page_no
            out.extend(_split(text, m))
    return out


def parse_ppt(path: str, base_meta: dict) -> list[tuple[str, dict]]:
    from pptx import Presentation
    out: list[tuple[str, dict]] = []
    prs = Presentation(path)
    for slide_no, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            elif getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    texts.append(" ".join(c.text for c in row.cells))
        text = "\n".join(t for t in texts if t.strip())
        if not text.strip():
            continue
        m = dict(base_meta); m["page"] = slide_no
        out.extend(_split(text, m))
    return out


def parse_word(path: str, base_meta: dict) -> list[tuple[str, dict]]:
    from docx import Document
    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    # 表格内容
    for tbl in doc.tables:
        for row in tbl.rows:
            text += "\n" + " ".join(c.text for c in row.cells)
    return _split(text, base_meta)


def parse_document(path: str, base_meta: dict) -> list[tuple[str, dict]]:
    """根据扩展名分派。"""
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return parse_pdf(path, base_meta)
    if ext in (".ppt", ".pptx"):
        return parse_ppt(path, base_meta)
    if ext in (".doc", ".docx"):
        return parse_word(path, base_meta)
    if ext == ".txt":
        return _split(Path(path).read_text(encoding="utf-8", errors="ignore"), base_meta)
    return []
