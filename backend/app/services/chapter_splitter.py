"""按章节标题切分 PDF 页面归属。

输入：PDF 全页文本 + 已知章节列表（来自 DB chapters 表）
输出：[{chapter_id, chapter_title, start_page, end_page}]

算法：
1. 跳过目录页（含点引导符 ... 或 ．．的页面）
2. 对每个非目录页，先匹配 "第N章" 模式；若无，匹配章节标题文本（如"数据类型与输入输出"）
3. 动态课程额外扫描页眉/页角（左上、右上顶栏），识别「第1章 引 言」等页眉标题
4. 每章在非目录页中的首次检测位置 = 该章起始页
5. 无匹配的页面回溯到最近检测到的章节（carry-forward）
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any

# 匹配 "第0章" / "第 1 章" / "第12章 概述" 等
_CHAPTER_PATTERN = re.compile(r"第\s*(\d+)\s*章")

# 点引导符（目录页特征）：英文省略号或中文全角点
_DOT_LEADER_PATTERN = re.compile(r"\.{3,}|．．+|…+")


def _is_toc_page(text: str) -> bool:
    """检测是否为目录页：含点引导符。"""
    if not text:
        return False
    # 只看前 800 字符（目录条目集中在页面前部）
    head = text[:800]
    return bool(_DOT_LEADER_PATTERN.search(head))


def _extract_title_text(full_title: str) -> str:
    """从 "第5章 数据类型与输入输出" 提取 "数据类型与输入输出"。"""
    m = _CHAPTER_PATTERN.match(full_title.strip())
    if m:
        rest = full_title.strip()[m.end():].strip()
        return rest
    return full_title.strip()


_CHAPTER_LINE_PATTERN = re.compile(
    r"第\s*(\d+)\s*章\s*(.+?)(?:\.{2,}|…{1,}|．．+|\s{3,}|$)"
)
_TOC_CHAPTER_LINE = re.compile(r"^第\s*(\d+)\s*章\s*(.+?)(?:\s+\d{1,4}\s*)?$")
_CHAPTER_HEADING_LINE = re.compile(r"^第\s*(\d+)\s*章\s+(.+)$")
_INVALID_SUBTITLE_PREFIXES = (
    "中讲过", "的先导", "已经", "，", "你", "在本书", "参见", "详见", "关于", "讲过",
)


def _clean_chapter_subtitle(subtitle: str) -> str:
    subtitle = subtitle.strip().strip(".·…|｜").strip()
    m = re.match(r"^(.+?)\s+\d{1,4}$", subtitle)
    if m:
        subtitle = m.group(1).strip()
    # 页眉常见「引 言」等汉字间空格
    subtitle = re.sub(r"(?<=[\u4e00-\u9fff])\s+(?=[\u4e00-\u9fff])", "", subtitle)
    return subtitle.strip()


def _normalize_header_line(text: str) -> str:
    """页眉行规范化：去竖线、首尾印刷页码等。"""
    s = (text or "").strip().strip("|｜").strip()
    s = re.sub(r"^\d{1,4}\s+", "", s)
    s = re.sub(r"\s+\d{1,4}\s*$", "", s)
    return s.strip()


def _is_english_book_title_line(line: str) -> bool:
    """页眉中的英文书名行（如 Database System Concepts, 6E）。"""
    line = (line or "").strip()
    if len(line) < 10 or _CHAPTER_PATTERN.search(line):
        return False
    ascii_chars = sum(1 for c in line if ord(c) < 128)
    return ascii_chars / max(len(line), 1) >= 0.75


def _parse_corner_header_block(corner_text: str) -> dict[str, Any] | None:
    """解析页眉/页角多行块（如 |第1章| + 英文书名 + 大标题「引 言」）。"""
    if not corner_text or not corner_text.strip():
        return None
    lines = [ln.strip() for ln in corner_text.splitlines() if ln.strip()]
    if not lines:
        return None

    for i, line in enumerate(lines[:8]):
        hit = parse_chapter_heading(_normalize_header_line(line))
        if not hit:
            continue
        if not hit.get("subtitle"):
            for nxt in lines[i + 1 : i + 4]:
                if _is_english_book_title_line(nxt):
                    continue
                cand = _clean_chapter_subtitle(nxt)
                if cand and _is_valid_chapter_subtitle(cand):
                    hit["subtitle"] = cand
                    hit["title"] = f"第{hit['order_idx']}章 {cand}"[:128]
                    break
        return hit

    for i in range(min(len(lines) - 1, 6)):
        combined = _normalize_header_line(f"{lines[i]} {lines[i + 1]}")
        hit = parse_chapter_heading(combined)
        if hit:
            return hit
    return None


def extract_pdf_pages_with_corners(path: str) -> tuple[list[str], list[str]]:
    """从 PDF 提取每页全文与页眉/页角区域文本（顶栏、左上、右上）。"""
    import pdfplumber

    pages_text: list[str] = []
    page_corners: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            full = page.extract_text() or ""
            pages_text.append(full)
            w, h = float(page.width or 0), float(page.height or 0)
            if w <= 0 or h <= 0:
                page_corners.append("")
                continue
            band_h = min(h * 0.18, 120)
            crops: list[str] = []
            for box in (
                (0, 0, w * 0.52, band_h),
                (w * 0.48, 0, w, band_h),
                (0, 0, w, band_h),
            ):
                try:
                    t = (page.crop(box).extract_text() or "").strip()
                except Exception:
                    t = ""
                if t:
                    crops.append(t)
            seen: set[str] = set()
            parts: list[str] = []
            for c in crops:
                if c not in seen:
                    seen.add(c)
                    parts.append(c)
            page_corners.append("\n".join(parts))
    return pages_text, page_corners


def _page_header_candidate_lines(full_text: str, corner_text: str = "") -> list[str]:
    """页眉/页角候选行：顶栏、页首/页末行（部分 PDF 右上角文本后渲染）。"""
    out: list[str] = []
    seen: set[str] = set()

    def add(s: str) -> None:
        s = _normalize_header_line(s)
        if not s or s in seen or len(s) > 120:
            return
        seen.add(s)
        out.append(s)

    block_hit = _parse_corner_header_block(corner_text)
    if block_hit:
        add(block_hit["title"])

    for src in (corner_text, full_text):
        if not src:
            continue
        lines = [ln.strip() for ln in src.splitlines() if ln.strip()]
        for ln in lines[:10]:
            add(ln)
        for i in range(min(len(lines) - 1, 8)):
            add(f"{lines[i]} {lines[i + 1]}")
        if len(lines) > 12:
            for ln in lines[-5:]:
                add(ln)
        add(src.strip()[:150].replace("\n", " "))
    return out


def _match_chapter_order_from_lines(
    lines: list[str],
    chapter_map: dict[int, dict[str, Any]],
    corner_text: str = "",
) -> int | None:
    """从候选行中匹配章节序号（优先 parse_chapter_heading，兼容原有行首匹配）。"""
    block = _parse_corner_header_block(corner_text)
    if block:
        order = block["order_idx"]
        subtitle = _clean_chapter_subtitle(block.get("subtitle") or "")
        if order in chapter_map and (not subtitle or _is_valid_chapter_subtitle(subtitle)):
            return order

    for line in lines:
        hit = parse_chapter_heading(line)
        if hit:
            order = hit["order_idx"]
            subtitle = _clean_chapter_subtitle(hit.get("subtitle") or "")
            if order in chapter_map and (not subtitle or _is_valid_chapter_subtitle(subtitle)):
                return order
        m = _CHAPTER_HEADING_LINE.match(line)
        if m:
            order = int(m.group(1))
            subtitle = _clean_chapter_subtitle(m.group(2))
            if order in chapter_map and _is_valid_chapter_subtitle(subtitle):
                return order
    return None


def _is_valid_chapter_subtitle(subtitle: str) -> bool:
    subtitle = subtitle.strip()
    if len(subtitle) < 2 or len(subtitle) > 80:
        return False
    if any(subtitle.startswith(p) for p in _INVALID_SUBTITLE_PREFIXES):
        return False
    if "，" in subtitle[:8]:
        return False
    return True


def _count_toc_chapter_lines(text: str) -> int:
    count = 0
    for line in text.splitlines()[:40]:
        line = line.strip()
        if _TOC_CHAPTER_LINE.match(line) or _CHAPTER_HEADING_LINE.match(line):
            count += 1
    return count


def _is_toc_page_for_extraction(text: str) -> bool:
    """目录页检测（用于动态课程；不因代码块中的 .....: 误判）。"""
    if not text:
        return False
    head = text[:1200].strip()
    if head.startswith("目录"):
        return True
    chapter_lines = _count_toc_chapter_lines(head)
    if chapter_lines >= 3:
        return True
    if chapter_lines >= 2 and _DOT_LEADER_PATTERN.search(head[:800]):
        return True
    return False


def split_pdf_by_chapters(
    pages_text: list[str],
    chapters: list[dict[str, Any]],
    *,
    dynamic_course: bool = False,
    page_corners: list[str] | None = None,
) -> list[dict[str, Any]]:
    """将 PDF 每页归到对应章节。

    Args:
        pages_text: 每页提取的文本，索引即页码(0-based)
        chapters: [{id, title, order_idx}] 来自 DB，按 order_idx 排序
        dynamic_course: True 时对 Java/Python 等动态课程使用增强目录/标题识别；
            False 时保持 C 语言等预置章节课程的原有拆分逻辑。

    Returns:
        [{chapter_id, chapter_title, order_idx, start_page, end_page}]
        start_page/end_page 为 1-based 页码，end_page 含本页
    """
    if not chapters or not pages_text:
        return []

    chapters_sorted = sorted(chapters, key=lambda c: c["order_idx"])
    chapter_map = {c["order_idx"]: c for c in chapters_sorted}

    def skip_toc_page(text: str) -> bool:
        return _is_toc_page_for_extraction(text) if dynamic_course else _is_toc_page(text)

    # 构造标题 → order 映射；短标题（如「数组」）仅允许整行标题匹配，避免正文误命中
    title_to_order: dict[str, int] = {}
    short_title_to_order: dict[str, int] = {}
    for c in chapters_sorted:
        title_text = _extract_title_text(c["title"])
        if not title_text:
            continue
        if len(title_text) >= 3:
            title_to_order[title_text] = c["order_idx"]
        elif len(title_text) >= 2:
            short_title_to_order[title_text] = c["order_idx"]

    def _normalize_heading_line(line: str) -> str:
        s = (line or "").strip()
        s = re.sub(r"^[「『\"'【\[]+", "", s)
        s = re.sub(r"[」』\"'】\]\s]+$", "", s)
        return s.strip()

    def _match_short_title_heading(text: str) -> int | None:
        """短章节名（2 字）仅当接近独立标题行时匹配，例如 OCR 成「数组」/「数组」」。"""
        if not short_title_to_order:
            return None
        for raw in (text or "").splitlines()[:5]:
            line = raw.strip()
            if not line or len(line) > 12:
                continue
            cleaned = _normalize_heading_line(line)
            order = short_title_to_order.get(cleaned)
            if order is None:
                continue
            # 行本身几乎就是标题（允许少量 OCR 标点）
            if len(_normalize_heading_line(line)) == len(cleaned) and len(line) <= len(cleaned) + 2:
                return order
        return None

    # 第一遍：为每个非目录页检测所属章节
    HEADER_LEN = 80
    page_chapter: dict[int, int] = {}  # page_idx(0-based) -> order_idx
    for i, text in enumerate(pages_text):
        if not text:
            continue
        if skip_toc_page(text):
            continue  # 跳过目录页

        if dynamic_course:
            for line in text.splitlines()[:12]:
                line = line.strip()
                if not line or len(line) > 80:
                    continue
                m = _CHAPTER_HEADING_LINE.match(line)
                if not m:
                    continue
                order = int(m.group(1))
                subtitle = _clean_chapter_subtitle(m.group(2))
                if order in chapter_map and _is_valid_chapter_subtitle(subtitle):
                    page_chapter[i] = order
                    break
            if i in page_chapter:
                continue
            # 补充：页眉/页角（左上、右上）中的「第N章 标题」
            corner = page_corners[i] if page_corners and i < len(page_corners) else ""
            order = _match_chapter_order_from_lines(
                _page_header_candidate_lines(text, corner),
                chapter_map,
                corner_text=corner,
            )
            if order is not None:
                page_chapter[i] = order
                continue

        head = text[:HEADER_LEN]
        # 1) 先匹配 "第N章" 模式（页眉）；动态课程仅接受行首章节标题，避免正文中「时间序列」等短语误匹配
        if not dynamic_course:
            m = _CHAPTER_PATTERN.search(head)
            if m:
                order = int(m.group(1))
                if order in chapter_map:
                    page_chapter[i] = order
                    continue

            # 2) 匹配章节标题文本（如 "数据类型与输入输出"）——仅限页眉区域
            matched = False
            for title_text, order in sorted(title_to_order.items(), key=lambda x: -len(x[0])):
                if title_text in head:
                    page_chapter[i] = order
                    matched = True
                    break
            if matched:
                continue

            # 3) 短标题独立行（「数组」被 OCR 丢了「第9章」时的常见情形）
            short_order = _match_short_title_heading(text)
            if short_order is not None:
                page_chapter[i] = short_order
                continue

    # 第二遍：carry-forward 填充无检测的页面
    # 单调推进：一旦进入第N章，不再回退到更早的章节（避免正文交叉引用"见第1章"误判）
    last_order: int | None = None
    for i in range(len(pages_text)):
        if skip_toc_page(pages_text[i] or ""):
            continue  # 目录页不分配
        if i in page_chapter:
            order = page_chapter[i]
            # 只接受 >= 当前章节的匹配（单调推进）
            if last_order is None or order >= last_order:
                last_order = order
            # 否则忽略这个匹配（视为交叉引用/页眉残留）
        if last_order is not None and i not in page_chapter:
            page_chapter[i] = last_order
        elif last_order is not None and page_chapter.get(i) != last_order:
            # 该页匹配的章节 < 当前章节，修正为当前章节
            page_chapter[i] = last_order

    # 第三遍：构建每章的页码区间
    if not page_chapter:
        # 全部未识别，归到第一章
        first = chapters_sorted[0]
        return [{
            "chapter_id": first["id"],
            "chapter_title": first["title"],
            "order_idx": first["order_idx"],
            "start_page": 1,
            "end_page": len(pages_text),
        }]

    # 按章节分组
    chapter_pages: dict[int, list[int]] = {}
    for page_idx, order in page_chapter.items():
        chapter_pages.setdefault(order, []).append(page_idx)

    # 构建结果，按 order_idx 排序
    result = []
    for order in sorted(chapter_pages.keys()):
        if order not in chapter_map:
            continue
        pages = sorted(chapter_pages[order])
        ch = chapter_map[order]
        result.append({
            "chapter_id": ch["id"],
            "chapter_title": ch["title"],
            "order_idx": order,
            "start_page": pages[0] + 1,  # 1-based
            "end_page": pages[-1] + 1,
            "page_count": len(pages),
        })

    # 处理第一章节起始页之前的非目录页（前言/序言）：归到第一章
    if result:
        first_start = result[0]["start_page"]
        if first_start > 1:
            # 找到 first_start 之前是否有非目录页
            has_preface = any(
                not skip_toc_page(pages_text[i] or "")
                for i in range(first_start - 1)
                if i < len(pages_text)
            )
            if has_preface:
                result[0]["start_page"] = 1

    return result


def detect_chapters_in_pdf(pages_text: list[str]) -> list[tuple[int, int]]:
    """检测 PDF 中出现的章节标题，返回 [(order_idx, page_number), ...]。
    用于单章上传时判断是否含多章内容。跳过目录页。
    """
    found = []
    seen = set()
    for page_idx, text in enumerate(pages_text):
        if not text or _is_toc_page(text):
            continue
        head = text[:250]
        m = _CHAPTER_PATTERN.search(head)
        if m:
            order = int(m.group(1))
            if order not in seen:
                seen.add(order)
                found.append((order, page_idx + 1))
    return found


def _add_definition(
    definitions: dict[int, dict[str, Any]],
    order: int,
    subtitle: str,
) -> None:
    subtitle = _clean_chapter_subtitle(subtitle)
    if subtitle and not _is_valid_chapter_subtitle(subtitle):
        return
    if order in definitions:
        if subtitle and not definitions[order].get("subtitle"):
            definitions[order]["subtitle"] = subtitle[:500]
            definitions[order]["title"] = f"第{order}章 {subtitle}"[:128]
        return
    title = f"第{order}章 {subtitle}".strip() if subtitle else f"第{order}章"
    definitions[order] = {
        "order_idx": order,
        "title": title[:128],
        "subtitle": subtitle[:500],
    }


def extract_chapter_definitions_from_pdf(
    pages_text: list[str],
    page_corners: list[str] | None = None,
    toc_page_start: int | None = None,
    toc_page_end: int | None = None,
) -> list[dict[str, Any]]:
    """从 PDF 目录页与正文章首页提取章节定义，供动态课程（Java/Python）创建 DB 章节。

    toc_page_start/end 为 1-based PDF 页码（与阅读器页码一致）；指定后优先从该区间提取目录。
    返回 [{order_idx, title, subtitle}]，title 形如「第1章 Java入门」。
    """
    definitions: dict[int, dict[str, Any]] = {}

    def _scan_toc_text(text: str) -> None:
        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue
            m = _CHAPTER_LINE_PATTERN.search(line)
            if m:
                _add_definition(definitions, int(m.group(1)), m.group(2))
                continue
            m = _TOC_CHAPTER_LINE.match(line)
            if m:
                _add_definition(definitions, int(m.group(1)), m.group(2))
                continue
            m = _CHAPTER_PATTERN.search(line)
            if m:
                _add_definition(definitions, int(m.group(1)), line[m.end():])

    # 1a) 用户指定的目录页（最高优先级）
    if toc_page_start is not None or toc_page_end is not None:
        start = max(1, int(toc_page_start or toc_page_end or 1))
        end = max(start, int(toc_page_end or toc_page_start or start))
        end = min(end, len(pages_text))
        for i in range(start - 1, end):
            text = pages_text[i] if i < len(pages_text) else ""
            if text:
                _scan_toc_text(text)

    # 1b) 自动检测目录页（指定页未识别出足够章节时补充）
    if len(definitions) < 2:
        for text in pages_text:
            if not text or not _is_toc_page_for_extraction(text):
                continue
            _scan_toc_text(text)

    # 2) 正文章首页 / 页眉：目录不足时从页眉块补充
    if len(definitions) < 2:
        for i, text in enumerate(pages_text):
            if not text or _is_toc_page_for_extraction(text):
                continue
            corner = page_corners[i] if page_corners and i < len(page_corners) else ""
            block = _parse_corner_header_block(corner)
            if block:
                _add_definition(definitions, block["order_idx"], block.get("subtitle") or "")
            for line in _page_header_candidate_lines(text, corner):
                hit = parse_chapter_heading(line)
                if hit:
                    _add_definition(definitions, hit["order_idx"], hit.get("subtitle") or "")
                    continue
                m = _CHAPTER_HEADING_LINE.match(line)
                if not m:
                    continue
                _add_definition(definitions, int(m.group(1)), m.group(2))

    return [definitions[k] for k in sorted(definitions.keys())]


_CN_DIGIT = {"零": 0, "〇": 0, "一": 1, "二": 2, "两": 2, "三": 3, "四": 4, "五": 5, "六": 6, "七": 7, "八": 8, "九": 9}


def parse_chinese_number(text: str) -> int | None:
    """解析中文或阿拉伯数字章节序号（1–99）。"""
    text = (text or "").strip()
    if not text:
        return None
    if text.isdigit():
        n = int(text)
        return n if 0 < n <= 99 else None
    if text == "十":
        return 10
    if text.startswith("十"):
        rest = text[1:]
        if not rest:
            return 10
        if rest in _CN_DIGIT:
            return 10 + _CN_DIGIT[rest]
        return None
    if text.endswith("十") and len(text) > 1:
        head = text[:-1]
        if head in _CN_DIGIT:
            return _CN_DIGIT[head] * 10
    if "十" in text:
        left, right = text.split("十", 1)
        tens = _CN_DIGIT.get(left, 1) if left else 1
        ones = _CN_DIGIT.get(right, 0) if right else 0
        n = tens * 10 + ones
        return n if 0 < n <= 99 else None
    if text in _CN_DIGIT:
        n = _CN_DIGIT[text]
        return n if n > 0 else None
    return None


_CHAPTER_ANY_PATTERN = re.compile(
    r"第\s*([一二三四五六七八九十两〇零\d]+)\s*章\s*(.*)$", re.I
)


def parse_chapter_heading(text: str) -> dict[str, Any] | None:
    """从单行文本解析章节（支持页眉「第 5 章 高级 SQL 93」「4 第 1 章 引 言」等）。"""
    text = _normalize_header_line(text)
    if not text or len(text) > 120:
        return None
    m = _CHAPTER_ANY_PATTERN.match(text)
    if not m:
        m = _CHAPTER_ANY_PATTERN.search(text)
    if not m:
        return None
    order = parse_chinese_number(m.group(1))
    if order is None or order <= 0 or order > 99:
        return None
    subtitle = _clean_chapter_subtitle(m.group(2))
    if subtitle and not _is_valid_chapter_subtitle(subtitle):
        if len(subtitle) > 80:
            subtitle = subtitle[:80]
        elif not subtitle.strip():
            subtitle = ""
        else:
            pass
    title = f"第{order}章 {subtitle}".strip() if subtitle else f"第{order}章"
    return {
        "order_idx": order,
        "title": title[:128],
        "subtitle": (subtitle or "")[:500],
    }


# 课件文件名：第1章 / 第七章 绪论.pptx / 02-变量.pdf
_FILENAME_PATTERNS: list[re.Pattern[str]] = [
    re.compile(r"^chapter\s*(\d+)\s*(.*)$", re.I),
    re.compile(r"^ch\s*(\d+)\s*(.*)$", re.I),
    re.compile(r"^(\d{1,2})[\.\-_、]\s*(.+)$"),
    re.compile(r"^(\d{1,2})\s+(.+)$"),
]


def extract_chapter_from_filename(filename: str) -> dict[str, Any] | None:
    """从课件文件名解析章节序号与标题。"""
    stem = Path(filename).stem.strip() if filename else ""
    if not stem:
        return None
    hit = parse_chapter_heading(stem)
    if hit:
        return hit
    for pat in _FILENAME_PATTERNS:
        m = pat.match(stem)
        if not m:
            continue
        order = int(m.group(1))
        if order <= 0 or order > 99:
            continue
        subtitle = (m.group(2) if m.lastindex and m.lastindex >= 2 else "").strip()
        subtitle = re.sub(r"[_\-]+$", "", subtitle).strip()
        title = f"第{order}章 {subtitle}".strip() if subtitle else f"第{order}章"
        return {
            "order_idx": order,
            "title": title[:128],
            "subtitle": subtitle[:500],
        }
    return None


def _lines_from_ppt(path: str) -> list[str]:
    from pptx import Presentation
    prs = Presentation(path)
    if not prs.slides:
        return []
    slide = prs.slides[0]
    lines: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    lines.append(t)
        elif getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                t = " ".join(c.text.strip() for c in row.cells if c.text.strip())
                if t:
                    lines.append(t)
    return lines


def _lines_from_word(path: str) -> list[str]:
    from docx import Document
    doc = Document(path)
    lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return lines[:20]


def try_extract_chapter_from_file(path: str) -> dict[str, Any] | None:
    """从课件内容（首页/首页幻灯片/文首）识别章节标题。"""
    ext = Path(path).suffix.lower()
    lines: list[str] = []
    try:
        if ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                if pdf.pages:
                    text = pdf.pages[0].extract_text() or ""
                    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        elif ext in (".ppt", ".pptx"):
            lines = _lines_from_ppt(path)
        elif ext in (".doc", ".docx"):
            lines = _lines_from_word(path)
        elif ext == ".txt":
            lines = Path(path).read_text(encoding="utf-8", errors="ignore").splitlines()[:20]
            lines = [ln.strip() for ln in lines if ln.strip()]
    except Exception:
        return None

    for line in lines[:20]:
        hit = parse_chapter_heading(line)
        if hit:
            return hit
    for i in range(min(len(lines) - 1, 5)):
        combined = f"{lines[i]} {lines[i + 1]}"
        hit = parse_chapter_heading(combined)
        if hit:
            return hit
    return None


def try_extract_chapter_from_pdf_first_page(path: str) -> dict[str, Any] | None:
    return try_extract_chapter_from_file(path)


def build_chapter_plan_from_courseware(
    filenames: list[str],
    *,
    file_paths: dict[int, str] | None = None,
) -> list[dict[str, Any]]:
    """根据课件文件名与内容生成章节计划。"""
    if not filenames:
        return []

    parsed: list[tuple[int, dict[str, Any], int, str]] = []
    unparsed: list[tuple[int, str]] = []

    for idx, name in enumerate(filenames):
        hit = extract_chapter_from_filename(name)
        source = "filename"
        if not hit and file_paths and idx in file_paths:
            hit = try_extract_chapter_from_file(file_paths[idx])
            source = "content"
        if hit:
            parsed.append((hit["order_idx"], hit, idx, source))
        else:
            unparsed.append((idx, name))

    if unparsed and file_paths:
        still: list[tuple[int, str]] = []
        for idx, name in unparsed:
            if idx in file_paths:
                hit = try_extract_chapter_from_file(file_paths[idx])
                if hit:
                    parsed.append((hit["order_idx"], hit, idx, "content"))
                    continue
            still.append((idx, name))
        unparsed = still

    used_orders = {item[0] for item in parsed}
    next_order = (max(used_orders) if used_orders else 0) + 1
    unparsed.sort(key=lambda x: x[1].lower())

    for idx, name in unparsed:
        while next_order in used_orders:
            next_order += 1
        stem = Path(name).stem.strip()
        hit = parse_chapter_heading(stem)
        if hit:
            parsed.append((hit["order_idx"], hit, idx, "filename"))
            used_orders.add(hit["order_idx"])
            continue
        if stem:
            title = f"第{next_order}章 {stem}"[:128]
        else:
            title = f"第{next_order}章"
        parsed.append((next_order, {"order_idx": next_order, "title": title, "subtitle": stem[:500]}, idx, "order"))
        used_orders.add(next_order)
        next_order += 1

    parsed.sort(key=lambda x: x[0])
    seen: set[int] = set()
    plan: list[dict[str, Any]] = []
    for order, defn, file_idx, source in parsed:
        if order in seen:
            raise ValueError(f"多个课件对应同一章节：第{order}章（请检查文件名或课件首页标题）")
        seen.add(order)
        plan.append({
            "order_idx": order,
            "title": defn["title"],
            "description": defn.get("subtitle") or "",
            "file_index": file_idx,
            "parse_source": source,
        })
    return plan
